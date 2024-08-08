from typing import List, Tuple, Dict, Any
# from io import BytesIO, StringIO
# from fastapi import UploadFile
# from pypdf import PdfReader
from urllib.parse import urlparse
from PIL import Image
# import urllib.request
import requests
import os
import pymupdf
import re
import pandas as pd
import pytesseract

from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_google_vertexai import VertexAIEmbeddings, VertexAI
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough, RunnableParallel
from langchain_core.output_parsers import JsonOutputParser,StrOutputParser
from langchain_core.pydantic_v1 import BaseModel, Field
from langchain.document_loaders import YoutubeLoader
from docx import Document as docu
from youtube_transcript_api import YouTubeTranscriptApi


from services.logger import setup_logger
from services.tool_registry import ToolFile
from api.error_utilities import LoaderError
from enum import Enum


#PowerPoint Loader imports
from pptx import Presentation
import os

from io import BytesIO
from langchain_core.documents import Document
from typing import List

#HTML and XML loaders
from bs4 import BeautifulSoup

#Extraction of all text from slides in presentation


relative_path = "features/quzzify"

logger = setup_logger(__name__)

def read_text_file(file_path):
    # Get the directory containing the script file
    script_dir = os.path.dirname(os.path.abspath(__file__))

    # Combine the script directory with the relative file path
    absolute_file_path = os.path.join(script_dir, file_path)
    
    with open(absolute_file_path, 'r') as file:
        return file.read()

class RAGRunnable:
    def __init__(self, func):
        self.func = func
    
    def __or__(self, other):
        def chained_func(*args, **kwargs):
            # Result of previous function is passed as first argument to next function
            return other(self.func(*args, **kwargs))
        return RAGRunnable(chained_func)
    
    def __call__(self, *args, **kwargs):
        return self.func(*args, **kwargs)

class YouTubeTranscriptLoader(BaseLoader):
    def __init__(self, verbose=False):
        self.verbose = verbose


    def fetch_transcript(self, video_id: str) -> str:
        try:
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            transcript_text = '\n'.join([transcript['text'] for transcript in transcript_list])
            return transcript_text
       
        except Exception as e:
            print(f"Error fetching transcript for video {video_id}: {e}")
            return ""


    def load(self, files: List[ToolFile]):
        documents = []


        for file in files:
            try:
                url = file.url
                video_id = YoutubeLoader.extract_video_id(url)
                transcript_text = self.fetch_transcript(video_id)
            
               
                if transcript_text:
                    documents.append(Document(page_content=transcript_text, metadata={'video_id': video_id}))
                    if self.verbose:
                        print(f"Fetched transcript for video {video_id}")


                else:
                    print(f"No transcript found for video {video_id}")


            except Exception as e:
                print(f"Error loading video {video_id}: {e}")
       
        return documents



class BytesFileCSVLoader(BaseLoader):
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() == "csv":
                # pdf_reader = PdfReader(file) #! PyPDF2.PdfReader is deprecated
                file.seek(0)
                df = pd.read_csv(file)
                for row in df.itertuples():
                    content = ""
                    for column in row[1:]:
                        content+= (str(column).strip() + "\n")
                    metadata = {"page_number": row[0] + 1, "source": file_type}
                    doc = Document(page_content=content, metadata=metadata)
                    documents.append(doc)               
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents

class BytesFileXLSXLoader(BaseLoader):

    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() == "xlsx":
                # pdf_reader = PdfReader(file) #! PyPDF2.PdfReader is deprecated
                file.seek(0)
                df = pd.read_excel(file)
                for row in df.itertuples():
                    content = ""
                    for column in row[1:]:
                        content+= (str(column).strip() + "\n")
                    metadata = {"page_number": row[0] + 1, "source": file_type}
                    doc = Document(page_content=content, metadata=metadata)
                    documents.append(doc)               
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents
     
class DocLoader(BaseLoader):

    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() == "docx":
                # pdf_reader = PdfReader(file) #! PyPDF2.PdfReader is deprecated
                docs = docu(file)
                for page_num, page in enumerate(docs.paragraphs):
                    page_content = ""
                    for paragraph in page.runs:
                        page_content += paragraph.text.strip() + "\n"
                    metadata = {"page_number": page_num + 1, "source": file_type}
                    doc = Document(page_content=page_content.rstrip(),metadata=metadata)
                    documents.append(doc)               
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents
    
class ImageLoader(BaseLoader):
    def __init__(self,files: List[Tuple[BytesIO,str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        text_completion_model = VertexAI(model='gemini-1.5-flash-001')
        prompt= PromptTemplate.from_template("I want you to check if there is any missing words in {text}. If there are any, I want to to autocomplete them with the most relevant word possible and make the whole thing grammatically correct. The output should be a string.")
        text_chain = (
                {"text": RunnablePassthrough()} 
                | prompt 
                | text_completion_model 
                | StrOutputParser()
            )

        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() in ['jpeg', 'jpg', 'png']:
                image = Image.open(file)
                text = pytesseract.image_to_string(image)
                result = text_chain.invoke({"text" : text})
                metadata = {"source":file_type,"page_number":1}
                document = Document(page_content=result,metadata=metadata)
                documents.append(document)
                    
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents
    

class BytesFilePDFLoader(BaseLoader):
    # Original def __init__(self, files: List[Tuple[BytesIO, str]])
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() == "pdf":
                pdf_reader = pymupdf.open(stream=file)
                for pages in range(pdf_reader.page_count):
                    page = pdf_reader.load_page(page_id=pages)
                    metadata = {"source" : file_type, "page_number" : pages + 1}
                    doc = Document(page_content=page.get_text(), metadata= metadata)
                    documents.append(doc)
                    
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents


class PowerPointLoader(BaseLoader):
    def __init__(self,files: List[Tuple[BytesIO, str]], loader = None, verbose=False, expected_file_type="pptx", ):
        self.loader = loader
        self.expected_file_type = expected_file_type
        self.verbose = verbose
        self.files = files
    
    def get_slide_text(slides):
        text_concepts = ""
        # Iterate over each shape in the slides collection
        for shape in slides.shapes:
            # Get the title of the slide
            title = ""
            if slides.shapes.title:
                title = slides.shapes.title.text
            texts = ""
            if shape.has_text_frame:
                # Extract text from each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Extract text from each run in the paragraph
                    for run in paragraph.runs:
                        texts += run.text
            text_concepts += texts
        return title, text_concepts

    def load(self) -> List[Document]:
        documents: List[Document] = []
        for file,file_type in self.files:
            if file_type not in ('pptx', 'ppt'):
                    raise ValueError(f"Unsupported file type: {file_type}")
            else:
                prs = Presentation(file)
                page_content = ""
                for slide_num, slide in enumerate(prs.slides, start = 1):
                    title, text_concepts = PowerPointLoader.get_slide_text(slide)
                    page_content += (title + text_concepts)
                    metadata = {"source": file_type, "page_number": slide_num}
                    doc = Document(page_content=page_content, metadata=metadata)
                    documents.append(doc)
        return documents
    
class HTMLLoader(BaseLoader):
    def __init__(self, files: List[Tuple[BytesIO, str]], expected_file_type="html", verbose=False):
        self.verbose = verbose
        self.expected_file_type = expected_file_type
        self.files = files

    def load(self) -> List[Document]:
        documents = []
        
        # Ensure file paths is a list
        for file, file_type in self.files:
            if file_type != "html":
                raise ValueError(f"Unsupported file type: {file_type}")
            else:
                byte_str = file.getvalue()
                # text_obj = byte_str.decode("utf-8")
                soup = BeautifulSoup(byte_str)
                text = soup.get_text()
                documents.append(Document(page_content=text, metadata={"source":file_type,"page_number":1}))

        return documents

class LocalFileLoader(BaseLoader):
    def __init__(self, file_paths: list[str], file_loader=None):
        self.file_paths = file_paths
        self.expected_file_types = ["xlsx", "pdf", "pptx", "csv", "docx", "jpeg", 'jpg', "png"]
        self.loader = file_loader or BytesFileXLSXLoader or BytesFilePDFLoader or BytesFileCSVLoader or DocLoader or ImageLoader
        self.loader_dict = {"xlsx":BytesFileXLSXLoader, "pdf":BytesFilePDFLoader, "pptx": PowerPointLoader, 
                        "csv": BytesFileCSVLoader, "docx": DocLoader,"jpeg": ImageLoader,
                        'jpg': ImageLoader,"png": ImageLoader, "ppt": PowerPointLoader, "html": HTMLLoader}
        
    def load(self) -> List[Document]:
        documents = []
        
        # Ensure file paths is a list
        self.file_paths = [self.file_paths] if isinstance(self.file_paths, str) else self.file_paths
    
        for file_path in self.file_paths:
            
            file_type = file_path.split(".")[-1]

            if file_type not in self.expected_file_types:
                exp_file_type = self.expected_file_types.join(", ")
                raise ValueError(f"Expected file types: {exp_file_type}, but got: {file_type}")

            with open(file_path, 'rb') as file:
                loader = self.loader_dict[file_type]
                documents.extend(loader([file]).load())
        return documents

class FileTypes(Enum):
    PDF = 'pdf'
    CSV = 'csv'
    PPTX = 'pptx'
    PPT = 'ppt'
    DOCX = "docx"
    XLSX = "xlsx"
    HTML = 'html'
    JPEG = 'jpeg'
    JPG = 'jpg'
    PNG = 'png'  
    
class URLLoader():
    def __init__(self, verbose=False):
        # self.expected_file_types = ["xlsx", "pdf", "pptx", "csv", "docx","jpeg",'jpg',"png", "ppt", "html"]
        self.verbose = verbose
        self.loader_dict = {"xlsx":BytesFileXLSXLoader, "pdf":BytesFilePDFLoader, "pptx": PowerPointLoader, 
                        "csv": BytesFileCSVLoader, "docx": DocLoader,"jpeg": ImageLoader,
                        'jpg': ImageLoader,"png": ImageLoader, "ppt": PowerPointLoader, "html": HTMLLoader}
        
    
    def download_from_drive(self,file_id : str):
        download_url = "https://docs.google.com/uc?export=download&id=" + file_id

        response = requests.get(download_url, stream=True)

        # Check for confirmation prompt
        if response.status_code == 302:  # Found a redirect, likely confirmation needed
            logger.info("Google Drive requires confirmation to download the file.")
            logger.info("Please visit the provided URL in your browser and allow access.")
            logger.info(response.headers['Location'])  # Print the redirection URL
            return None  # Indicate download not completed
        
        # Download logic (assuming confirmation was successful)
        file_type = ''
        content_disposition = response.headers.get('Content-Disposition')
        if content_disposition:
            filename_part = content_disposition.split('=')[-1]
            if '.' in filename_part:
                file_type = filename_part.split('.')[-1].lower()[:len(filename_part.split('.')[-1]) - 1]

        return (response,file_type)
    
    def load(self, tool_files: List[ToolFile]) -> List[Document]:
        queued_files = []
        documents = []
        youtube_files = []
        # any_success = False
        response = None

        def check_file_type(file_type):
            return any(file_type == member.value for member in FileTypes)

        for tool_file in tool_files:
            url = tool_file.url
            file_type = None
            regex = r"/d/([^?]+)/"
            
            if url.lower().startswith("https://youtu.be/"):
                youtube_files.append(tool_file)
                                           
            try:
                match = re.search(regex,url)
                if not match:
                    response = requests.get(url, verify=False, stream=True)
                    parsed_url = urlparse(url)
                    path = parsed_url.path
                else:
                    
                    file_id = match.group(1)
                    response,file_type = self.download_from_drive(file_id)
                if response.status_code == 200:
                    # Read file
                    file_content = BytesIO(response.content)
                    # Check file type
                    # file_type = path.rsplit(".")[-1]
                    if not file_type:
                        file_type = url.rsplit('.')[-1]
                    # if file_type not in self.expected_file_types:
                    if not check_file_type(file_type):
                        # string = self.expected_file_types.join(", ")
                        filelist = [str(member.value) for member in FileTypes]
                        string = ''.join(filelist)
                        raise LoaderError(f"Expected file types: {string}, but got: {file_type}")

                    # Append to Queue
                    queued_files.append((file_content, file_type))
                    if self.verbose:
                        logger.info(f"Successfully loaded file from {url}")
                    
                else:
                    logger.error(f"Request failed to load file from {url} and got status code {response.status_code}")

            except Exception as e:
                logger.error(f"Failed to load file from {url}")
                logger.error(e)
                continue
        if youtube_files:
            yt_loader = YouTubeTranscriptLoader(verbose=self.verbose)
            docs = yt_loader.load(youtube_files)
            if self.verbose:
                print(f"Documents from YouTube loader: {len(docs)}")
                documents.extend(docs)
            if self.verbose:
                print(f"Total documents: {len(documents)}")
                return documents   
     
        # Pass Queue to the file loader if there are any successful loads
        if len(queued_files) > 0:
            documents = []
            for file in queued_files: # run each file one by one
                loader = self.loader_dict[file[1]]
                # loader = self.loader_dict[FileTypes.file[1]]
                file_loader = loader([file])
                try:
                    documents.extend(file_loader.load())
                    if self.verbose:
                        logger.info(f"Loaded {len(documents)} documents")
                except: # some error
                    continue
        else:
            raise LoaderError("Unable to load any files from URLs")

        return documents
  
class RAGpipeline:
    def __init__(self, loader=None, splitter=None, vectorstore_class=None, embedding_model=None, verbose=False):
        default_config = {
            "loader": URLLoader(verbose = verbose),
            "splitter": RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100),
            "vectorstore_class": Chroma,
            "embedding_model": VertexAIEmbeddings(model='textembedding-gecko')
        }
        self.loader = loader or default_config["loader"]
        self.splitter = splitter or default_config["splitter"]
        self.vectorstore_class = vectorstore_class or default_config["vectorstore_class"]
        self.embedding_model = embedding_model or default_config["embedding_model"]
        self.verbose = verbose

    def load_PDFs(self, files) -> List[Document]:
        if self.verbose:
            logger.info(f"Loading {len(files)} files")
            logger.info(f"Loader type used: {type(self.loader)}")
        
        logger.debug(f"Loader is a: {type(self.loader)}")
        
        try:
            total_loaded_files = self.loader.load(files)
        except LoaderError as e:
            logger.error(f"Loader experienced error: {e}")
            raise LoaderError(e)
            
        return total_loaded_files
    
    def split_loaded_documents(self, loaded_documents: List[Document]) -> List[Document]:
        if self.verbose:
            logger.info(f"Splitting {len(loaded_documents)} documents")
            logger.info(f"Splitter type used: {type(self.splitter)}")
            
        total_chunks = []
        chunks = self.splitter.split_documents(loaded_documents)
        total_chunks.extend(chunks)
        
        if self.verbose: logger.info(f"Split {len(loaded_documents)} documents into {len(total_chunks)} chunks")
        
        return total_chunks

    def create_vectorstore(self, documents: List[Document]):
        if self.verbose:
            logger.info(f"Creating vectorstore from {len(documents)} documents")
            for document in documents:
                logger.info(document)
        try:
            self.vectorstore = self.vectorstore_class.from_documents(documents, self.embedding_model)
            logger.info(f"Vectorstore created")
        except Exception as e:
            logger.error(f"Error creating vectorstore: {e}")
            raise  # Rethrow the exception to handle it further
        
        if self.verbose:
            logger.info(f"Vectorstore created")
        
        return self.vectorstore
    
    def compile(self):
        # Compile the pipeline
        self.load_PDFs = RAGRunnable(self.load_PDFs)
        logger.info("Completed loading PDFs")
        self.split_loaded_documents = RAGRunnable(self.split_loaded_documents)
        logger.info("Completed splitting loaded documents")
        self.create_vectorstore = RAGRunnable(self.create_vectorstore)
        if self.verbose: logger.info(f"Completed pipeline compilation")
    
    def __call__(self, documents):
        # Returns a vectorstore ready for usage 
        
        if self.verbose: 
            logger.info(f"Executing pipeline")
            logger.info(f"Start of Pipeline received: {len(documents)} documents of type {type(documents[0])}")
        
        pipeline = self.load_PDFs | self.split_loaded_documents | self.create_vectorstore
        return pipeline(documents)

class QuizBuilder:
    def __init__(self, vectorstore, topic, prompt=None, model=None, parser=None, verbose=False):
        default_config = {
            "model": VertexAI(model="gemini-1.0-pro"), 
            "parser": JsonOutputParser(pydantic_object=QuizQuestion),
            "prompt": read_text_file("prompt/quizzify-prompt.txt")
        }
        
        self.prompt = prompt or default_config["prompt"]
        self.model = model or default_config["model"]
        self.parser = parser or default_config["parser"]
        
        self.vectorstore = vectorstore
        self.topic = topic
        self.verbose = verbose
        
        if vectorstore is None: raise ValueError("Vectorstore must be provided")
        if topic is None: raise ValueError("Topic must be provided")
    
    def compile(self):
        # Return the chain
        prompt = PromptTemplate(
            template=self.prompt,
            input_variables=["topic"],
            partial_variables={"format_instructions": self.parser.get_format_instructions()}
        )
        
        retriever = self.vectorstore.as_retriever()
        
        runner = RunnableParallel(
            {"context": retriever, "topic": RunnablePassthrough()}
        )
        
        chain = runner | prompt | self.model | self.parser
        
        if self.verbose: logger.info(f"Chain compilation complete")
        
        return chain

    def validate_response(self, response: Dict) -> bool:
        try:
            # Assuming the response is already a dictionary
            if isinstance(response, dict):
                if 'question' in response and 'choices' in response and 'answer' in response and 'explanation' in response:
                    choices = response['choices']
                    if isinstance(choices, dict):
                        for key, value in choices.items():
                            if not isinstance(key, str) or not isinstance(value, str):
                                return False
                        return True
            return False
        except TypeError as e:
            if self.verbose:
                logger.error(f"TypeError during response validation: {e}")
            return False

    def format_choices(self, choices: Dict[str, str]) -> List[Dict[str, str]]:
        return [{"key": k, "value": v} for k, v in choices.items()]
    
    def create_questions(self, num_questions: int = 5) -> List[Dict]:
        if self.verbose: logger.info(f"Creating {num_questions} questions")
        
        if num_questions > 10:
            return {"message": "error", "data": "Number of questions cannot exceed 10"}
        
        chain = self.compile()
        
        generated_questions = []
        attempts = 0
        max_attempts = num_questions * 10  # Allow for more attempts to generate questions

        while len(generated_questions) < num_questions and attempts < max_attempts:
            response = chain.invoke(self.topic)
            if self.verbose:
                logger.info(f"Generated response attempt {attempts + 1}: {response}")
            
            # Directly check if the response format is valid
            if self.validate_response(response):
                response["choices"] = self.format_choices(response["choices"])
                generated_questions.append(response)
                if self.verbose:
                    logger.info(f"Valid question added: {response}")
                    logger.info(f"Total generated questions: {len(generated_questions)}")
            else:
                if self.verbose:
                    logger.warning(f"Invalid response format. Attempt {attempts + 1} of {max_attempts}")
            
            # Move to the next attempt regardless of success to ensure progress
            attempts += 1

        # Log if fewer questions are generated
        if len(generated_questions) < num_questions:
            logger.warning(f"Only generated {len(generated_questions)} out of {num_questions} requested questions")
        
        if self.verbose: logger.info(f"Deleting vectorstore")
        self.vectorstore.delete_collection()
        
        # Return the list of questions
        return generated_questions[:num_questions]

class QuestionChoice(BaseModel):
    key: str = Field(description="A unique identifier for the choice using letters A, B, C, D, etc.")
    value: str = Field(description="The text content of the choice")
class QuizQuestion(BaseModel):
    question: str = Field(description="The question text")
    choices: List[QuestionChoice] = Field(description="A list of choices")
    answer: str = Field(description="The correct answer")
    explanation: str = Field(description="An explanation of why the answer is correct")
